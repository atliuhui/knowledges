"""Document conversion: original formats -> AI-friendly processed/.

Each converter returns (text, target_suffix). Conversion is best-effort; missing
optional dependencies for binary formats degrade to a clear error rather than a crash.
"""

from __future__ import annotations

from dataclasses import dataclass
import html
import json
from pathlib import Path
from typing import Any, Callable


@dataclass
class ConversionResult:
    text: str
    target_suffix: str  # e.g. ".md" or ".csv"


class ConversionError(RuntimeError):
    pass


def _normalize_text_content(text: str) -> str:
    # Normalize line endings and always end with a single trailing newline.
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    return text.rstrip("\n") + "\n"


def _normalize_csv_content(text: str) -> str:
    # Keep raw CSV content, but normalize BOM/newlines for stable downstream ingest.
    text = text.lstrip("\ufeff")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    return text.rstrip("\n") + "\n"


def _convert_text(path: Path) -> ConversionResult:
    text = path.read_text(encoding="utf-8", errors="replace")
    return ConversionResult(_normalize_text_content(text), path.suffix.lower())


def _convert_plain(path: Path) -> ConversionResult:
    text = path.read_text(encoding="utf-8", errors="replace")
    return ConversionResult(_normalize_text_content(text), path.suffix.lower())


def _convert_docx(path: Path) -> ConversionResult:
    try:
        from docx import Document  # type: ignore
    except ImportError as e:
        raise ConversionError(f"python-docx not installed: {e}") from e
    doc = Document(str(path))
    lines: list[str] = []
    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        if style.startswith("heading"):
            # "heading 1" -> level 1
            try:
                level = int(style.split()[-1])
            except (ValueError, IndexError):
                level = 2
            lines.append(f"{'#' * max(1, min(level, 6))} {text}")
        else:
            lines.append(text)
    # Tables
    for table in doc.tables:
        lines.append("")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    return ConversionResult("\n".join(lines).strip() + "\n", ".md")


def _convert_pdf(path: Path) -> ConversionResult:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError as e:
        raise ConversionError(f"pypdf not installed: {e}") from e
    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = page.extract_text() or ""
        except Exception as e:  # noqa: BLE001
            txt = f"<!-- page {i} extract failed: {e} -->"
        parts.append(f"## Page {i}\n\n{txt.strip()}")
    return ConversionResult("\n\n".join(parts).strip() + "\n", ".md")


def _convert_pptx(path: Path) -> ConversionResult:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise ConversionError(f"python-pptx not installed: {e}") from e
    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {idx}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"\n> Notes: {notes}")
        parts.append("")
    return ConversionResult("\n".join(parts).strip() + "\n", ".md")


def _convert_xlsx(path: Path) -> ConversionResult:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as e:
        raise ConversionError(f"openpyxl not installed: {e}") from e
    import csv
    import io

    wb = load_workbook(str(path), data_only=True, read_only=True)
    out = io.StringIO()
    writer = csv.writer(out)
    for sheet in wb.worksheets:
        writer.writerow([f"# sheet: {sheet.title}"])
        for row in sheet.iter_rows(values_only=True):
            writer.writerow(["" if v is None else v for v in row])
        writer.writerow([])
    return ConversionResult(out.getvalue(), ".csv")


def _convert_csv(path: Path) -> ConversionResult:
    text = path.read_text(encoding="utf-8-sig", errors="replace")
    return ConversionResult(_normalize_csv_content(text), ".csv")


def _convert_drawio(path: Path) -> ConversionResult:
    import base64
    import re
    import zlib
    from urllib.parse import unquote
    from xml.etree import ElementTree as ET

    def _strip_html(text: str) -> str:
        text = html.unescape(text)
        text = text.replace("<br>", "\n").replace("<br/>", "\n").replace("<br />", "\n")
        if "<" not in text:
            return text
        # draw.io labels are usually tiny HTML fragments; strip tags while keeping visible text.
        return re.sub(r"<[^>]+>", "", text)

    def _diagram_xml(diagram_el: ET.Element) -> str:
        if len(diagram_el):
            # Uncompressed format: <diagram><mxGraphModel>...</mxGraphModel></diagram>
            return ET.tostring(diagram_el[0], encoding="unicode")
        raw = (diagram_el.text or "").strip()
        if not raw:
            return ""
        if raw.startswith("<"):
            return raw
        try:
            decoded = base64.b64decode(raw)
            inflated = zlib.decompress(decoded, -15)
            return unquote(inflated.decode("utf-8", errors="replace"))
        except Exception as e:  # noqa: BLE001
            raise ConversionError(f"failed to decode drawio compressed diagram: {e}") from e

    def _extract_values(diagram_xml: str) -> list[str]:
        if not diagram_xml.strip():
            return []
        try:
            root = ET.fromstring(diagram_xml)
        except Exception as e:  # noqa: BLE001
            raise ConversionError(f"failed to parse drawio diagram xml: {e}") from e
        values: list[str] = []
        for cell in root.findall(".//mxCell"):
            v = (cell.attrib.get("value") or "").strip()
            if v:
                try:
                    clean = _strip_html(v).strip()
                except Exception:
                    clean = html.unescape(v)
                if clean:
                    values.append(clean)
        return values

    try:
        root = ET.parse(str(path)).getroot()
    except Exception as e:  # noqa: BLE001
        raise ConversionError(f"failed to parse drawio file: {e}") from e

    parts: list[str] = []
    diagrams = root.findall(".//diagram")
    if not diagrams and root.tag == "diagram":
        diagrams = [root]
    for idx, d in enumerate(diagrams, start=1):
        name = (d.attrib.get("name") or f"Diagram {idx}").strip()
        parts.append(f"## Diagram {idx}: {name}")
        vals = _extract_values(_diagram_xml(d))
        if vals:
            parts.extend(f"- {v}" for v in vals)
        else:
            parts.append("(no text labels)")
        parts.append("")

    if not parts:
        raise ConversionError("drawio file contains no diagram")
    return ConversionResult("\n".join(parts).strip() + "\n", ".md")


def _convert_xmind(path: Path) -> ConversionResult:
    import zipfile

    def _walk_topic(topic: dict, lines: list[str], level: int = 1) -> None:
        title = str(topic.get("title") or "").strip()
        if title:
            lines.append(f"{'#' * min(level, 6)} {title}")

        notes = topic.get("notes") or {}
        if isinstance(notes, dict):
            plain = notes.get("plain")
            if isinstance(plain, dict):
                content = str(plain.get("content") or "").strip()
                if content:
                    lines.append(content)

        children = topic.get("children") or {}
        attached = children.get("attached") if isinstance(children, dict) else None
        if isinstance(attached, list):
            for child in attached:
                if isinstance(child, dict):
                    _walk_topic(child, lines, level + 1)

    if not zipfile.is_zipfile(path):
        raise ConversionError("xmind file is not a valid zip archive")

    with zipfile.ZipFile(path, "r") as zf:
        names = set(zf.namelist())
        if "content.json" in names:
            try:
                data = json.loads(zf.read("content.json").decode("utf-8", errors="replace"))
            except Exception as e:  # noqa: BLE001
                raise ConversionError(f"failed to parse xmind content.json: {e}") from e

            sheets = data if isinstance(data, list) else [data]
            lines: list[str] = []
            for idx, sheet in enumerate(sheets, start=1):
                if not isinstance(sheet, dict):
                    continue
                title = str(sheet.get("title") or f"Sheet {idx}").strip()
                lines.append(f"# Sheet {idx}: {title}")
                root_topic = sheet.get("rootTopic")
                if isinstance(root_topic, dict):
                    _walk_topic(root_topic, lines, level=2)
                lines.append("")
            if not lines:
                raise ConversionError("xmind content.json has no readable sheets")
            return ConversionResult("\n".join(lines).strip() + "\n", ".md")

        if "content.xml" in names:
            from xml.etree import ElementTree as ET

            ns = {
                "xmap": "urn:xmind:xmap:xmlns:content:2.0",
                "fo": "http://www.w3.org/1999/XSL/Format",
            }

            try:
                root = ET.fromstring(zf.read("content.xml"))
            except Exception as e:  # noqa: BLE001
                raise ConversionError(f"failed to parse xmind content.xml: {e}") from e

            lines: list[str] = []
            sheets = root.findall("xmap:sheet", ns)
            for idx, sheet in enumerate(sheets, start=1):
                s_title = sheet.findtext("xmap:title", default=f"Sheet {idx}", namespaces=ns).strip()
                lines.append(f"# Sheet {idx}: {s_title}")
                for topic in sheet.findall(".//xmap:topic", ns):
                    title = topic.findtext("xmap:title", default="", namespaces=ns).strip()
                    if title:
                        lines.append(f"- {title}")
                lines.append("")

            if not lines:
                raise ConversionError("xmind content.xml has no readable topics")
            return ConversionResult("\n".join(lines).strip() + "\n", ".md")

    raise ConversionError("unsupported xmind structure: missing content.json/content.xml")


_RAPID_OCR_CACHE: dict[tuple[str, str, str, str], Any] = {}


def _get_rapid_ocr(version: str, model_type: str, lang_det: str, lang_rec: str):
    """Lazily build and cache a RapidOCR engine per (version, model_type, det lang, rec lang)."""
    key = (version, model_type, lang_det, lang_rec)
    cached = _RAPID_OCR_CACHE.get(key)
    if cached is not None:
        return cached

    rapid_cls = None
    enums_mod = None
    try:
        from rapidocr import RapidOCR as _RapidOCRNew  # type: ignore
        rapid_cls = _RapidOCRNew
        try:
            import rapidocr as enums_mod  # type: ignore
        except ImportError:
            enums_mod = None
    except ImportError:
        try:
            from rapidocr_onnxruntime import RapidOCR as _RapidOCRLegacy  # type: ignore
            rapid_cls = _RapidOCRLegacy
        except ImportError as e:
            raise ConversionError(
                "rapidocr not installed: reinstall the project with `pip install -e .`"
            ) from e

    def _resolve_enum(enum_name: str, value: str):
        """Map a string like 'PP-OCRv5' / 'mobile' / 'ch' to the matching RapidOCR enum value.

        rapidocr 2.x requires Enum-typed values in the params dict; we look the
        enum up by name (case-insensitive, ignoring punctuation) and fall back
        to a member whose .value matches the requested string.
        """
        if enums_mod is None:
            return value
        enum_cls = getattr(enums_mod, enum_name, None)
        if enum_cls is None:
            return value
        norm = value.strip().upper().replace("-", "").replace("_", "")
        for member in enum_cls:  # type: ignore[attr-defined]
            mname = member.name.upper().replace("-", "").replace("_", "")
            if mname == norm:
                return member
            mval = str(getattr(member, "value", "")).upper().replace("-", "").replace("_", "")
            if mval == norm:
                return member
        return value  # let RapidOCR raise a clearer error if truly invalid

    # rapidocr 2.x supports per-stage model selection via a params dict; values
    # must be Enum-typed (OCRVersion / ModelType / LangDet / LangRec).
    params = {
        "Det.ocr_version": _resolve_enum("OCRVersion", version),
        "Det.model_type": _resolve_enum("ModelType", model_type),
        "Det.lang_type": _resolve_enum("LangDet", lang_det),
        "Rec.ocr_version": _resolve_enum("OCRVersion", version),
        "Rec.model_type": _resolve_enum("ModelType", model_type),
        "Rec.lang_type": _resolve_enum("LangRec", lang_rec),
    }
    attempts: list[dict[str, Any]] = [{"params": params}, {}]
    last_err: Exception | None = None
    ocr = None
    for kwargs in attempts:
        try:
            ocr = rapid_cls(**kwargs)
            break
        except Exception as e:  # noqa: BLE001
            last_err = e
            continue
    if ocr is None:
        raise ConversionError(f"failed to initialize RapidOCR: {last_err}")

    _RAPID_OCR_CACHE[key] = ocr
    return ocr


def _rapid_extract_texts(raw_result: Any) -> list[str]:
    """Normalize RapidOCR output across 1.x / 2.x return shapes into a text list."""
    lines: list[str] = []
    if raw_result is None:
        return lines

    # rapidocr 2.x: a RapidOCRResult object with .txts (list[str]).
    if hasattr(raw_result, "txts"):
        for t in (getattr(raw_result, "txts", None) or []):
            s = str(t).strip()
            if s:
                lines.append(s)
        return lines

    # rapidocr 1.x / rapidocr_onnxruntime: (result, elapse) tuple where
    # result is a list of [box, text, confidence] entries (or None when empty).
    if isinstance(raw_result, tuple):
        raw_result = raw_result[0]
    if not raw_result:
        return lines
    for entry in raw_result:
        try:
            text = entry[1]
        except (IndexError, TypeError):
            continue
        if isinstance(text, (list, tuple)) and text:
            text = text[0]
        s = str(text).strip() if text is not None else ""
        if s:
            lines.append(s)
    return lines


def _convert_image(path: Path, options: dict[str, Any] | None = None) -> ConversionResult:
    opts = options or {}
    version = str(opts.get("image_ocr_version") or "PP-OCRv5").strip() or "PP-OCRv5"
    model_type = str(opts.get("image_ocr_model_type") or "mobile").strip() or "mobile"
    lang_det = str(opts.get("image_ocr_lang_det") or "ch").strip() or "ch"
    lang_rec = str(opts.get("image_ocr_lang_rec") or "ch").strip() or "ch"
    ocr = _get_rapid_ocr(version, model_type, lang_det, lang_rec)
    try:
        raw = ocr(str(path))
    except Exception as e:  # noqa: BLE001
        raise ConversionError(f"rapidocr failed on {path.name}: {e}") from e

    lines = _rapid_extract_texts(raw)
    header = f"# {path.stem}"
    if not lines:
        body = f"<!-- no text recognized in {path.name} -->"
    else:
        body = "\n".join(lines)
    return ConversionResult(f"{header}\n\n{body}\n", ".md")


_REGISTRY: dict[str, Callable[[Path], ConversionResult]] = {
    ".md": _convert_text,
    ".markdown": _convert_text,
    ".txt": _convert_plain,
    ".rst": _convert_plain,
    ".yaml": _convert_plain,
    ".yml": _convert_plain,
    ".json": _convert_plain,
    ".docx": _convert_docx,
    ".pdf": _convert_pdf,
    ".pptx": _convert_pptx,
    ".xlsx": _convert_xlsx,
    ".csv": _convert_csv,
    ".tsv": _convert_csv,
    ".drawio": _convert_drawio,
    ".xmind": _convert_xmind,
}

# Image formats are handled separately because they need converter options (lang).
_IMAGE_SUFFIXES: frozenset[str] = frozenset({".png", ".jpg", ".jpeg"})


def convert(path: Path, options: dict[str, Any] | None = None) -> ConversionResult:
    """Convert a source file to AI-friendly text (Markdown / CSV)."""
    suffix = path.suffix.lower()
    if suffix in _IMAGE_SUFFIXES:
        return _convert_image(path, options)
    handler = _REGISTRY.get(suffix)
    if not handler:
        raise ConversionError(f"unsupported format: {suffix}")
    return handler(path)
